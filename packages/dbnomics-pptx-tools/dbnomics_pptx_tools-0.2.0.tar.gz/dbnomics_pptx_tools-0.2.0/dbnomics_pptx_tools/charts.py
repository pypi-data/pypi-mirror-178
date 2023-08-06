from typing import Iterable, cast

import pandas as pd
from pandas import DatetimeIndex
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.chart.datalabel import DataLabel
from pptx.chart.point import Point
from pptx.chart.series import _BaseCategorySeries
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import Slide

from dbnomics_pptx_tools.formatters import format_number
from dbnomics_pptx_tools.metadata import ChartSpec, DataLabelPosition
from dbnomics_pptx_tools.pptx_copy import copy_shape_properties
from dbnomics_pptx_tools.repo import SeriesRepo

__all__ = ["update_chart"]


# Cf https://github.com/scanny/python-pptx/issues/716
def add_data_label_to_last_point_of_each_series(
    chart: Chart,
    *,
    chart_spec: ChartSpec,
    repo: SeriesRepo,
):
    for series_index, series in enumerate(cast(Iterable[_BaseCategorySeries], chart.series)):
        points = list(cast(Iterable[Point], series.points))
        last_point = points[-1]
        data_label = cast(DataLabel, last_point.data_label)
        data_label.has_text_frame = True
        data_label.text_frame.text = build_data_label(series.name, chart_spec=chart_spec, repo=repo)
        data_label.position = pick_data_label_position(series_index)


def build_data_label(
    series_name: str,
    *,
    chart_spec: ChartSpec,
    repo: SeriesRepo,
) -> str:
    """Build a data label for a single point.

    DataLabel for single Point does not implement show_series_name, show_category and show_value,
    like DataLabel does, so fallback to manual solution.
    """
    series_id = chart_spec.find_series_id_by_name(series_name)
    if series_id is None:
        raise ValueError(f"Could not find the series ID from the name {series_name!r}")
    df = repo.load(series_id)
    latest_observation = cast(float, df.sort_values(by="period").iloc[-1].value)
    return f"{series_name}: {format_number(latest_observation)}"


def build_category_chart_data(chart_spec: ChartSpec, *, repo: SeriesRepo) -> CategoryChartData:
    chart_spec_series_ids = chart_spec.get_series_ids()
    chart_data = CategoryChartData()
    df = pd.concat(repo.load(series_id) for series_id in chart_spec_series_ids)
    pivoted_df = df.pivot(index="period", columns="series_id", values="value")
    chart_data.categories = cast(DatetimeIndex, pivoted_df.index).to_pydatetime()
    for series_id in chart_spec_series_ids:
        series_spec = chart_spec.find_series_spec(series_id)
        if series_spec is None:
            raise ValueError(f"Could not find spec for series {series_id!r}")
        series_name = series_spec.name
        series = pivoted_df[series_id].fillna("")
        chart_data.add_series(series_name, series.values)
    return chart_data


def pick_data_label_position(series_index: int) -> XL_DATA_LABEL_POSITION:
    """Pick a position for a data label.

    When many series points are close, data labels can overlap.
    The DataLabel implementation does not seem to have a "no overlap" option.
    This is a workaround that works with 2 points, for now, due to the "modulo" mechanics.
    """
    candidates = [
        XL_DATA_LABEL_POSITION.ABOVE,  # type: ignore
        XL_DATA_LABEL_POSITION.BELOW,  # type: ignore
    ]
    return candidates[series_index % len(candidates)]


def update_chart(
    chart_shape: GraphicFrame,
    *,
    chart_spec: ChartSpec,
    repo: SeriesRepo,
    slide: Slide,
):
    chart = cast(Chart, chart_shape.chart)
    chart_data = build_category_chart_data(chart_spec, repo=repo)
    chart_shape.element.getparent().remove(chart_shape.element)
    new_chart_shape = cast(
        GraphicFrame,
        cast(SlideShapes, slide.shapes).add_chart(
            chart.chart_type, chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height, chart_data
        ),
    )
    copy_shape_properties(chart_shape, new_chart_shape)
    new_chart = cast(Chart, new_chart_shape.chart)
    if DataLabelPosition.LAST_POINT.value in chart_spec.data_labels:
        add_data_label_to_last_point_of_each_series(new_chart, chart_spec=chart_spec, repo=repo)
